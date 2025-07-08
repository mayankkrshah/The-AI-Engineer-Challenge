"""
Multi-format file loader supporting comprehensive file type processing.
Supports: PDF, DOCX, XLSX, PPTX, TXT, HTML, MD, JSON, YAML, CSV, RTF, and code files.
"""

import os
import mimetypes
from typing import List, Dict, Optional, Tuple
from pathlib import Path
import logging

# Try to import magic (optional dependency)
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False

# Document processing imports
try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

try:
    import json
    JSON_AVAILABLE = True
except ImportError:
    JSON_AVAILABLE = False

try:
    import csv
    CSV_AVAILABLE = True
except ImportError:
    CSV_AVAILABLE = False

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class MultiFormatLoader:
    """
    Comprehensive file loader supporting multiple formats for RAG applications.
    """
    
    SUPPORTED_FORMATS = {
        # Document formats (modern only)
        'pdf': 'PDF documents',
        'docx': 'Word documents (modern format only)',
        'rtf': 'Rich Text Format',
        'txt': 'Plain text files',
        'md': 'Markdown files',
        'html': 'HTML files',
        'htm': 'HTML files',
        
        # Spreadsheet formats (modern only)
        'xlsx': 'Excel spreadsheets (modern format only)',
        'csv': 'CSV files',
        
        # Presentation formats (modern only)
        'pptx': 'PowerPoint presentations (modern format only)',
        
        # Data formats
        'json': 'JSON files',
        'yaml': 'YAML files',
        'yml': 'YAML files',
        'xml': 'XML files',
        
        # Code formats (web3 and general)
        'sol': 'Solidity smart contracts',
        'py': 'Python files',
        'js': 'JavaScript files',
        'ts': 'TypeScript files',
        'jsx': 'React JSX files',
        'tsx': 'React TSX files',
        'rs': 'Rust files',
        'go': 'Go files',
        'css': 'CSS files',
        'scss': 'SCSS files',
        'sass': 'Sass files',
        'less': 'Less files'
    }
    
    # Legacy formats that are NOT supported
    UNSUPPORTED_LEGACY_FORMATS = {
        'doc': 'Legacy Word documents (.doc) are not supported. Please convert to .docx format.',
        'ppt': 'Legacy PowerPoint presentations (.ppt) are not supported. Please convert to .pptx format.',
        'xls': 'Legacy Excel spreadsheets (.xls) are not supported. Please convert to .xlsx format.'
    }
    
    def __init__(self, file_path: str):
        """
        Initialize the multi-format loader.
        
        Args:
            file_path: Path to the file to be loaded
        """
        self.file_path = file_path
        self.file_extension = Path(file_path).suffix.lower().lstrip('.')
        self.mime_type = self._detect_mime_type()
        
        logger.info(f"MultiFormatLoader initialized for: {file_path}")
        logger.info(f"Detected extension: {self.file_extension}")
        logger.info(f"Detected MIME type: {self.mime_type}")
    
    def _detect_mime_type(self) -> str:
        """Detect file MIME type using multiple methods."""
        mime_type = None
        
        # Try python-magic first (most accurate) if available
        if MAGIC_AVAILABLE:
            try:
                mime_type = magic.from_file(self.file_path, mime=True)
            except Exception as e:
                logger.warning(f"python-magic failed: {e}")
        
        # Fallback to mimetypes module
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(self.file_path)
        
        return mime_type or 'application/octet-stream'
    
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding for text files."""
        if not CHARDET_AVAILABLE:
            return 'utf-8'
        
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # Read first 10KB
                result = chardet.detect(raw_data)
                return result['encoding'] or 'utf-8'
        except Exception as e:
            logger.warning(f"Encoding detection failed: {e}")
            return 'utf-8'
    
    def is_supported(self) -> bool:
        """Check if the file format is supported."""
        # Check for legacy formats first
        if self.file_extension in self.UNSUPPORTED_LEGACY_FORMATS:
            return False
        return self.file_extension in self.SUPPORTED_FORMATS
    
    def get_file_info(self) -> Dict[str, str]:
        """Get information about the file."""
        # Check if it's a legacy format first
        if self.file_extension in self.UNSUPPORTED_LEGACY_FORMATS:
            description = self.UNSUPPORTED_LEGACY_FORMATS[self.file_extension]
        else:
            description = self.SUPPORTED_FORMATS.get(self.file_extension, 'Unknown format')
        
        return {
            'filename': os.path.basename(self.file_path),
            'extension': self.file_extension,
            'mime_type': self.mime_type,
            'supported': self.is_supported(),
            'description': description
        }
    
    def load(self) -> List[str]:
        """
        Load and extract text content from the file.
        
        Returns:
            List of text chunks extracted from the file
        """
        # Check for legacy formats first
        if self.file_extension in self.UNSUPPORTED_LEGACY_FORMATS:
            error_msg = self.UNSUPPORTED_LEGACY_FORMATS[self.file_extension]
            raise ValueError(error_msg)
        
        if not self.is_supported():
            supported_formats = ', '.join(sorted(self.SUPPORTED_FORMATS.keys()))
            raise ValueError(f"Unsupported file format: {self.file_extension}. Supported formats: {supported_formats}")
        
        # Check if file exists and is not empty
        if not os.path.exists(self.file_path):
            raise FileNotFoundError(f"File not found: {self.file_path}")
        
        if os.path.getsize(self.file_path) == 0:
            raise ValueError("The uploaded file is empty")
        
        try:
            # Handle different file types
            if self.file_extension == 'pdf':
                return self._load_pdf()
            elif self.file_extension == 'docx':
                return self._load_docx()
            elif self.file_extension == 'pptx':
                return self._load_powerpoint()
            elif self.file_extension == 'xlsx':
                return self._load_excel()
            elif self.file_extension == 'csv':
                return self._load_csv()
            elif self.file_extension in ['json']:
                return self._load_json()
            elif self.file_extension in ['yaml', 'yml']:
                return self._load_yaml()
            elif self.file_extension == 'xml':
                return self._load_xml()
            elif self.file_extension == 'rtf':
                return self._load_rtf()
            elif self.file_extension in ['html', 'htm']:
                return self._load_html()
            elif self.file_extension in ['txt', 'md', 'py', 'js', 'ts', 'jsx', 'tsx', 'rs', 'go', 'css', 'scss', 'sass', 'less', 'sol']:
                return self._load_text()
            else:
                raise ValueError(f"Unsupported file format: {self.file_extension}")
        
        except Exception as e:
            logger.error(f"Error loading file {self.file_path}: {str(e)}")
            # Provide user-friendly error messages
            if "Package not found" in str(e):
                raise ValueError(f"The {self.file_extension.upper()} file appears to be corrupted or invalid. Please try uploading a different file.")
            elif "codec can't decode" in str(e):
                raise ValueError("Unable to read file - the file may be corrupted or in an unsupported text encoding.")
            else:
                raise ValueError(f"Unable to process the {self.file_extension.upper()} file. Please check that the file is valid and not corrupted.")
    
    def _load_pdf(self) -> List[str]:
        """Load PDF files."""
        if not PYPDF2_AVAILABLE:
            raise ValueError("PDF files are not supported in this installation. Please contact support for PDF processing capabilities.")
        
        try:
            reader = PdfReader(self.file_path)
            text_chunks = []
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_chunks.append(f"Page {page_num}:\n{text}")
            
            if not text_chunks:
                raise ValueError("No readable text could be extracted from this PDF file. The PDF may be image-based or corrupted.")
            
            return text_chunks
        except Exception as e:
            if "PDF files are not supported" in str(e):
                raise
            logger.error(f"Error reading PDF: {e}")
            raise ValueError(f"Unable to process this PDF file. The file may be corrupted, password-protected, or in an unsupported format.")
    
    def _load_docx(self) -> List[str]:
        """Load Word documents."""
        if not DOCX_AVAILABLE:
            raise ValueError("Word documents (.docx/.doc) are not supported in this installation. Please contact support for Word document processing capabilities.")
        
        try:
            doc = Document(self.file_path)
            text_chunks = []
            current_chunk = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    current_chunk += paragraph.text + "\n"
                    
                    # Create chunk every 5 paragraphs or when encountering headings
                    if (len(current_chunk) > 1000 or 
                        any(run.bold for run in paragraph.runs) or
                        paragraph.style.name.startswith('Heading')):
                        if current_chunk.strip():
                            text_chunks.append(current_chunk.strip())
                            current_chunk = ""
            
            # Add remaining content
            if current_chunk.strip():
                text_chunks.append(current_chunk.strip())
            
            if not text_chunks:
                raise ValueError("No readable text could be extracted from this Word document. The document may be empty or corrupted.")
            
            return text_chunks
        except Exception as e:
            if "Word documents" in str(e):
                raise
            logger.error(f"Error reading DOCX: {e}")
            raise ValueError(f"Unable to process this Word document. The file may be corrupted or in an unsupported format.")
    
    def _load_excel(self) -> List[str]:
        """Load Excel files."""
        if not EXCEL_AVAILABLE:
            raise ValueError("Excel files (.xlsx/.xls) are not supported in this installation. Please contact support for Excel processing capabilities.")
        
        try:
            workbook = load_workbook(self.file_path, read_only=True)
            text_chunks = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = f"Sheet: {sheet_name}\n\n"
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                        sheet_content += row_text + "\n"
                
                if sheet_content.strip():
                    text_chunks.append(sheet_content)
            
            if not text_chunks:
                raise ValueError("No readable data could be extracted from this Excel file. The spreadsheet may be empty or corrupted.")
            
            return text_chunks
        except Exception as e:
            if "Excel files" in str(e):
                raise
            logger.error(f"Error reading Excel: {e}")
            raise ValueError(f"Unable to process this Excel file. The file may be corrupted, password-protected, or in an unsupported format.")
    
    def _load_powerpoint(self) -> List[str]:
        """Load PowerPoint presentations."""
        if not PPTX_AVAILABLE:
            raise ValueError("PowerPoint files (.pptx/.ppt) are not supported in this installation. Please contact support for PowerPoint processing capabilities.")
        
        try:
            prs = Presentation(self.file_path)
            text_chunks = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"Slide {slide_num}:\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content += shape.text + "\n"
                
                if slide_content.strip():
                    text_chunks.append(slide_content)
            
            if not text_chunks:
                raise ValueError("No readable text could be extracted from this PowerPoint presentation. The slides may be empty or contain only images.")
            
            return text_chunks
        except Exception as e:
            if "PowerPoint files" in str(e):
                raise
            logger.error(f"Error reading PowerPoint: {e}")
            raise ValueError(f"Unable to process this PowerPoint file. The file may be corrupted or in an unsupported format.")
    
    def _load_rtf(self) -> List[str]:
        """Load RTF files."""
        if not RTF_AVAILABLE:
            raise ValueError("RTF files are not supported in this installation. Please contact support for RTF processing capabilities.")
        
        try:
            with open(self.file_path, 'r', encoding='utf-8') as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                
                if not text.strip():
                    raise ValueError("No readable text could be extracted from this RTF file. The document may be empty or corrupted.")
                
                return [text]
        except Exception as e:
            if "RTF files are not supported" in str(e):
                raise
            logger.error(f"Error reading RTF: {e}")
            raise ValueError(f"Unable to process this RTF file. The file may be corrupted or in an unsupported format.")
    
    def _load_html(self) -> List[str]:
        """Load HTML files."""
        if not HTML_AVAILABLE:
            raise ImportError("beautifulsoup4 is required for HTML processing")
        
        try:
            encoding = self._detect_encoding(self.file_path)
            with open(self.file_path, 'r', encoding=encoding) as file:
                html_content = file.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Remove script and style elements
                for element in soup(["script", "style"]):
                    element.decompose()
                
                text = soup.get_text()
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Error reading HTML: {e}")
            raise
    
    def _load_markdown(self) -> List[str]:
        """Load Markdown files."""
        try:
            encoding = self._detect_encoding(self.file_path)
            with open(self.file_path, 'r', encoding=encoding) as file:
                md_content = file.read()
                
                if MARKDOWN_AVAILABLE:
                    # Convert markdown to HTML then to text
                    html = markdown.markdown(md_content)
                    if HTML_AVAILABLE:
                        soup = BeautifulSoup(html, 'html.parser')
                        text = soup.get_text()
                    else:
                        text = md_content
                else:
                    text = md_content
                
                return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Error reading Markdown: {e}")
            raise
    
    def _load_json(self) -> List[str]:
        """Load JSON files."""
        if not JSON_AVAILABLE:
            raise ImportError("json is required for JSON processing")
        
        try:
            encoding = self._detect_encoding(self.file_path)
            with open(self.file_path, 'r', encoding=encoding) as file:
                data = json.load(file)
                
                # Convert JSON to readable text
                if isinstance(data, dict):
                    text = self._dict_to_text(data)
                elif isinstance(data, list):
                    text = self._list_to_text(data)
                else:
                    text = str(data)
                
                return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Error reading JSON: {e}")
            raise
    
    def _load_yaml(self) -> List[str]:
        """Load YAML files."""
        if not YAML_AVAILABLE:
            raise ImportError("PyYAML is required for YAML processing")
        
        try:
            encoding = self._detect_encoding(self.file_path)
            with open(self.file_path, 'r', encoding=encoding) as file:
                data = yaml.safe_load(file)
                
                # Convert YAML to readable text
                if isinstance(data, dict):
                    text = self._dict_to_text(data)
                elif isinstance(data, list):
                    text = self._list_to_text(data)
                else:
                    text = str(data)
                
                return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Error reading YAML: {e}")
            raise
    
    def _load_csv(self) -> List[str]:
        """Load CSV files."""
        if not CSV_AVAILABLE:
            raise ImportError("csv is required for CSV processing")
        
        try:
            encoding = self._detect_encoding(self.file_path)
            text_chunks = []
            
            with open(self.file_path, 'r', encoding=encoding) as file:
                csv_reader = csv.reader(file)
                headers = next(csv_reader, None)
                
                if headers:
                    text_chunks.append(f"CSV Headers: {', '.join(headers)}")
                
                chunk_content = ""
                for row_num, row in enumerate(csv_reader, 1):
                    if row:
                        chunk_content += f"Row {row_num}: {', '.join(row)}\n"
                        
                        # Create chunk every 100 rows
                        if row_num % 100 == 0:
                            text_chunks.append(chunk_content)
                            chunk_content = ""
                
                # Add remaining content
                if chunk_content.strip():
                    text_chunks.append(chunk_content)
            
            return text_chunks
        except Exception as e:
            logger.error(f"Error reading CSV: {e}")
            raise
    
    def _load_xml(self) -> List[str]:
        """Load XML files."""
        if not HTML_AVAILABLE:
            raise ImportError("beautifulsoup4 is required for XML processing")
        
        try:
            encoding = self._detect_encoding(self.file_path)
            with open(self.file_path, 'r', encoding=encoding) as file:
                xml_content = file.read()
                soup = BeautifulSoup(xml_content, 'xml')
                text = soup.get_text()
                
                return [text] if text.strip() else []
        except Exception as e:
            logger.error(f"Error reading XML: {e}")
            raise
    
    def _load_text(self) -> List[str]:
        """Load plain text files (including code files)."""
        encodings_to_try = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'iso-8859-1']
        
        # Start with detected encoding if available
        detected_encoding = self._detect_encoding(self.file_path)
        if detected_encoding and detected_encoding not in encodings_to_try:
            encodings_to_try.insert(0, detected_encoding)
        
        content = None
        last_error = None
        
        for encoding in encodings_to_try:
            try:
                with open(self.file_path, 'r', encoding=encoding, errors='ignore') as file:
                    content = file.read()
                    break
            except Exception as e:
                last_error = e
                continue
        
        if content is None:
            logger.error(f"Could not read text file with any encoding: {last_error}")
            raise ValueError(f"Unable to read file - the file may be corrupted or in an unsupported text encoding")
        
        if not content.strip():
            raise ValueError("The file appears to be empty or contains no readable text")
        
        # For code files, add context about the file type
        if self.file_extension in ['sol', 'js', 'ts', 'py', 'rs', 'go', 'css', 'scss', 'sass', 'less']:
            content = f"File type: {self.SUPPORTED_FORMATS[self.file_extension]}\n\n{content}"
        
        return [content]
    
    def _dict_to_text(self, data: dict, indent: int = 0) -> str:
        """Convert dictionary to readable text."""
        lines = []
        prefix = "  " * indent
        
        for key, value in data.items():
            if isinstance(value, dict):
                lines.append(f"{prefix}{key}:")
                lines.append(self._dict_to_text(value, indent + 1))
            elif isinstance(value, list):
                lines.append(f"{prefix}{key}:")
                lines.append(self._list_to_text(value, indent + 1))
            else:
                lines.append(f"{prefix}{key}: {value}")
        
        return "\n".join(lines)
    
    def _list_to_text(self, data: list, indent: int = 0) -> str:
        """Convert list to readable text."""
        lines = []
        prefix = "  " * indent
        
        for i, item in enumerate(data):
            if isinstance(item, dict):
                lines.append(f"{prefix}[{i}]:")
                lines.append(self._dict_to_text(item, indent + 1))
            elif isinstance(item, list):
                lines.append(f"{prefix}[{i}]:")
                lines.append(self._list_to_text(item, indent + 1))
            else:
                lines.append(f"{prefix}[{i}]: {item}")
        
        return "\n".join(lines)
    
    @classmethod
    def get_supported_formats(cls) -> Dict[str, str]:
        """Get all supported file formats."""
        return cls.SUPPORTED_FORMATS.copy()
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Get list of supported file extensions."""
        return list(cls.SUPPORTED_FORMATS.keys()) 